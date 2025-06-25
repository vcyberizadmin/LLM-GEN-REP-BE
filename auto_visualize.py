#!/usr/bin/env python3
"""auto_visualize.py

Automate creation of basic visualizations from a zipped bundle of slide folders.

Usage:
    python auto_visualize.py path/to/bundle.zip [--pptx]

Dependencies:
    - pandas
    - matplotlib
    - (optional) python-pptx for --pptx output
"""

import argparse
import csv
import sys
import zipfile
from pathlib import Path
from tempfile import TemporaryDirectory
import re
import warnings

import pandas as pd
import matplotlib.pyplot as plt


# Mapping of keyword sets to chart types
_CHART_RULES = {
    frozenset(["trend", "timeseries", "growth"]): "line",
    frozenset(["distribution", "hist", "density"]): "hist",
    frozenset(["share", "comparison", "rank", "top"]): "bar",
    frozenset(["correlation", "relationship", "scatter"]): "scatter",
}


def choose_chart_type(title: str) -> str:
    """Return chart type based on keywords found in the title."""
    title_l = title.lower()
    for keywords, chart in _CHART_RULES.items():
        if any(k in title_l for k in keywords):
            return chart
    return "bar"


def safe_filename(text: str) -> str:
    """Sanitize text for file names."""
    return re.sub(r"[^A-Za-z0-9_.-]+", "_", text).strip("_")


def load_and_merge(folder: Path):
    """Load CSV/XLSX files from folder and merge on common columns."""
    dataframes = []
    names = []
    for f in folder.iterdir():
        if f.suffix.lower() == ".csv":
            df = pd.read_csv(f)
        elif f.suffix.lower() in {".xlsx", ".xls", ".xlsm", ".xlsb"}:
            df = pd.read_excel(f)
        else:
            continue
        dataframes.append(df)
        names.append(f.name)
    if not dataframes:
        raise ValueError("No supported data files found")
    merged = dataframes[0]
    for df in dataframes[1:]:
        common = [c for c in merged.columns if c in df.columns]
        if common:
            merged = merged.merge(df, on=common, how="left")
        else:
            # No common columns; keep the first dataset only
            break
    return merged, names, len(merged)


def plot_dataframe(df: pd.DataFrame, chart_type: str, title: str, out_file: Path):
    """Generate and save plot based on chart type."""
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    if not numeric_cols:
        warnings.warn(f"Slide '{title}' has no numeric columns; skipping plot")
        return False

    plt.clf()
    fig, ax = plt.subplots()

    if chart_type == "hist":
        ax.hist(df[numeric_cols[0]].dropna(), bins=20)
    elif chart_type == "scatter" and len(numeric_cols) >= 2:
        ax.scatter(df[numeric_cols[0]], df[numeric_cols[1]])
    elif chart_type == "line" and len(numeric_cols) >= 1:
        x = df.index
        if len(numeric_cols) > 1:
            y = numeric_cols[1]
            x = df[numeric_cols[0]]
        else:
            y = numeric_cols[0]
        ax.plot(x, df[y])
    else:  # bar chart or fallback
        cat_col = df.columns[0]
        y = numeric_cols[0]
        ax.bar(df[cat_col], df[y])
    ax.set_title(title)
    fig.autofmt_xdate()
    fig.tight_layout()
    fig.savefig(out_file)
    plt.close(fig)
    return True


def assemble_pptx(images: list[Path], out_file: Path):
    """Create PowerPoint from images if python-pptx is available."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        sys.stderr.write("python-pptx not installed; skipping PPTX creation\n")
        return

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for img in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img), Inches(0.5), Inches(0.5), width=prs.slide_width - Inches(1))
    prs.save(out_file)


def process_zip(zip_path: Path, create_ppt: bool = False, output_dir: Path | None = None):
    """Process a zipped slide bundle and return slide information.

    Parameters
    ----------
    zip_path : Path
        Path to the input ZIP file.
    create_ppt : bool, optional
        Whether to assemble the resulting images into a PowerPoint file.
    output_dir : Path, optional
        Directory where images (and PPTX) are written. Defaults to a
        sibling ``output`` folder next to this script.

    Returns
    -------
    tuple[list[dict], Path | None]
        A list of slide info dictionaries and the PowerPoint path if
        created.
    """

    out_dir = output_dir or (Path(__file__).resolve().parent / "output")
    out_dir.mkdir(exist_ok=True)
    generated_images = []
    slide_infos = []
    pptx_path: Path | None = None

    with TemporaryDirectory() as tmpdir:
        with zipfile.ZipFile(zip_path) as zf:
            zf.extractall(tmpdir)
        tmp_path = Path(tmpdir)
        # find slide folders
        slide_dirs = []
        for p in tmp_path.iterdir():
            if p.is_dir():
                m = re.match(r"Slide-(\d+)-(.*)", p.name)
                if m:
                    slide_dirs.append((int(m.group(1)), m.group(2), p))
        slide_dirs.sort(key=lambda x: x[0])

        for num, title, folder in slide_dirs:
            try:
                df, names, rows = load_and_merge(folder)
                chart_type = choose_chart_type(title)
                img_name = f"Slide-{num}-{safe_filename(title)}.png"
                img_path = out_dir / img_name
                plotted = plot_dataframe(df, chart_type, title, img_path)
                if not plotted:
                    continue
                generated_images.append(img_path)
                main_file = names[0]
                slide_infos.append({
                    "slide_num": num,
                    "title": title,
                    "chart_type": chart_type,
                    "main_file": main_file,
                    "rows": rows,
                    "image": img_path,
                })
                print(
                    f"[\u2713] Slide {num} – '{title}' → {chart_type} using {main_file} ({rows} rows)"
                )
            except Exception as e:
                print(f"[!] Slide {num} – '{title}' skipped: {e}", file=sys.stderr)

    if create_ppt and generated_images:
        pptx_path = out_dir / "visuals.pptx"
        assemble_pptx(generated_images, pptx_path)

    return slide_infos, pptx_path


def main():
    parser = argparse.ArgumentParser(description="Automate slide visualizations")
    parser.add_argument("zipfile", type=Path, help="Path to ZIP bundle")
    parser.add_argument("--pptx", action="store_true", help="Assemble PNGs into PowerPoint if python-pptx is installed")
    args = parser.parse_args()

    if not args.zipfile.exists():
        parser.error(f"ZIP file not found: {args.zipfile}")

    slides, pptx_path = process_zip(args.zipfile, args.pptx)
    for s in slides:
        print(
            f"[\u2713] Slide {s['slide_num']} – '{s['title']}' → {s['chart_type']} using {s['main_file']} ({s['rows']} rows)"
        )
    if pptx_path:
        print(f"PowerPoint created at: {pptx_path}")


if __name__ == "__main__":
    main()
